# First, ensure you have the library installed:
# pip install python-pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# --- New, Concise Presentation Content ---
slides_data = [
    {
        "title": "SAOCOM InSAR DEM Validation",
        "content": "A Land Cover-Stratified Analysis\n\nKey Pillars:\n• Vertical Accuracy Assessment\n• Land Cover Performance\n• Spatial Coverage & Void Analysis"
    },
    {
        "title": "Datasets & Study Area",
        "content": "SAOCOM: L-band InSAR point cloud\n\nReference Data:\n• TINITALY 10m DEM\n• Copernicus 30m DEM\n• CORINE 2018 Land Cover\n\nStudy Area: Verona, Italy\n• Constrained to SAOCOM data hull\n• All data projected to UTM Zone 32N"
    },
    {
        "title": "Core Processing Parameters",
        "content": "Coherence Filter: γ ≥ 0.3\n• Purpose: Removes noisy, unstable points to ensure data quality.\n\nGrid Resolution: 10 meters\n• Purpose: Matches the highest-resolution reference (TINITALY).\n\nResampling:\n• Cubic Convolution for elevation (smooths surfaces).\n• Nearest Neighbor for land cover (preserves classes)."
    },
    {
        "title": "Data Preprocessing Workflow",
        "content": "SAOCOM:\n1. Filter by coherence (γ ≥ 0.3).\n2. Remove invalid points.\n3. Apply k-NN algorithm to remove spatial outliers.\n\nReference DEMs:\n1. Reproject to UTM 32N.\n2. Resample to 10m grid.\n3. Clip to SAOCOM data boundary."
    },
    {
        "title": "Height Calibration",
        "content": "Goal: Convert relative SAOCOM heights to an absolute reference.\n\nMethod: Calculated the median offset using high-coherence points (γ ≥ 0.8).\n\nResult (vs. TINITALY):\n• Offset: +4.308 m\n• Validation RMSE: 4.955 m at stable points."
    },
    {
        "title": "Reference DEM Cross-Comparison",
        "content": "Purpose: Establish a baseline of agreement between TINITALY and Copernicus.\n\nKey Metrics:\n• RMSE: 4.68 m\n• Mean Bias: -2.03 m (TINITALY lower)\n• Correlation: 0.999\n\nTakeaway: The reference datasets themselves are not identical. This difference provides context for the SAOCOM validation."
    },
    {
        "title": "Statistical Metrics Used",
        "content": "Classical Metrics:\n• Bias (ME): Systematic error.\n• RMSE: Overall error magnitude (sensitive to outliers).\n\nRobust Metrics:\n• Median: Central tendency (resists outliers).\n• NMAD: Robust standard deviation.\n\nWhy robust? Our data has outliers; these metrics give a truer picture of performance."
    },
    {
        "title": "Height Statistics Summary",
        "content": "This table shows the fundamental statistical distributions for all elevation datasets. Note the large range of the raw SAOCOM data and the differences in mean/median values, underscoring the need for calibration.\n\n[Placeholder for Table 1: Height Statistics Summary]"
    },
    {
        "title": "Reference DEM Comparison",
        "content": "This figure visually breaks down the differences between the two reference DEMs. The maps highlight spatial patterns of disagreement, while the histogram shows the statistical distribution.\n\n[Placeholder for Figure 1: 8-Panel Reference DEM Comparison]"
    },
    {
        "title": "SAOCOM Spatial Coverage Analysis",
        "content": "A grid-based analysis reveals the extent of data gaps.\n\nVoid Percentage: 87.0%\n\nThis means for every 10 pixels in the study area, nearly 9 have no SAOCOM data.\n\nKey Question: What land cover types are in these voids?"
    },
    {
        "title": "CORINE Land Cover",
        "content": "The CORINE 2018 dataset was processed to classify the terrain.\n\n• 10 unique classes were found in the study area.\n• A custom color palette was used for clarity in all subsequent maps and charts.\n\n[Placeholder for Figure 4: CORINE Land Cover Map of Study Area]"
    },
    {
        "title": "Linking Elevation to Land Cover",
        "content": "Each SAOCOM point was tagged with the specific land cover class it falls on.\n\nThis creates the final analysis dataset, enabling us to ask:\n• \"How does accuracy change between forests and urban areas?\"\n• \"Which land cover types have the most data voids?\""
    },
    {
        "title": "Height Residuals by Land Cover",
        "content": "This table shows the error statistics for each land cover class. Performance varies significantly. Urban areas show lower error (NMAD < 3.5m), while forests show much higher error (NMAD > 5m).\n\n[Placeholder for Table 2: Height Residuals by Land Cover]"
    },
    {
        "title": "Void Analysis by Land Cover",
        "content": "This table quantifies which land cover classes are most affected by data gaps. Water bodies and forests are almost entirely voids, while vineyards and forests are the largest contributors to the total void area.\n\n[Placeholder for Table 3: Void Analysis by Land Cover]"
    },
    {
        "title": "Error vs. Coherence",
        "content": "These plots show the clear relationship between signal quality (coherence) and accuracy. As coherence increases, the error distribution tightens significantly, and the median error approaches zero.\n\n[Placeholder for Figure 2: Violin Plots of Residuals by Coherence Bins]"
    },
    {
        "title": "Error by Land Cover",
        "content": "These plots visualize the error distributions for each major land cover class. Note the wide shape for forests (high uncertainty) compared to the narrow shape for urban fabric (more reliable).\n\n[Placeholder for Figure 3: Side-by-Side Violin Plots of Residuals by Land Cover]"
    },
    {
        "title": "Individual Land Cover Overlay",
        "content": "To understand the spatial distribution, each class was mapped over a satellite image. This example shows the extensive \"Vineyards\" class, a key agricultural feature in the region.\n\n[Placeholder for Figure 5: Example Land Cover Map (e.g., Vineyards)]"
    },
    {
        "title": "Gridded SAOCOM Residuals",
        "content": "The point-based errors were interpolated into a continuous map. This allows us to visualize the spatial patterns of where the SAOCOM DEM is higher or lower than the reference DEMs.\n\n[Placeholder for Figure 6: 6-Panel Gridded SAOCOM Residual Map]"
    },
    {
        "title": "Height Correlation",
        "content": "These plots directly compare elevation values. While there is a strong linear relationship (high correlation), the spread of points around the 1:1 line visually represents the error.\n\n[Placeholder for Figure 7: Height Comparison Scatter Plots]"
    },
    {
        "title": "Land Cover Inside the Voids",
        "content": "This map shows only the land cover present in areas where SAOCOM data is missing. It clearly demonstrates that forests, vegetation, and water dominate the data gaps, pointing to decorrelation as the cause.\n\n[Placeholder for Figure 8: Map of Land Cover within Void Zones]"
    },
    {
        "title": "Quantifying Void Contributors",
        "content": "These charts pinpoint the most problematic land cover classes. The left chart shows the worst relative coverage (water), while the right chart shows the largest contributors to total void area (forests).\n\n[Placeholder for Figure 9: Bar Charts of Void Analysis by Land Cover]"
    },
    {
        "title": "Coverage and Voids Across the Landscape",
        "content": "This map shows all land cover types, with white areas indicating the SAOCOM data voids. This single image provides a powerful summary of the 87% void statistic and the fragmented nature of the coverage.\n\n[Placeholder for Figure 10: Land Cover Map with SAOCOM Void Zones]"
    },
    {
        "title": "Coverage vs. Voids: A Closer Look",
        "content": "This map zooms in on a single class, partitioning it into areas of successful SAOCOM coverage versus void areas. This is essential for assessing the data's utility for specific applications (e.g., forestry).\n\n[Placeholder for Figure 11: Example Coverage/Void Map (e.g., Broad-leaved forest)]"
    },
    {
        "title": "Conclusions",
        "content": "Accuracy: SAOCOM can produce accurate elevation data, but performance is highly dependent on coherence and land cover type.\n\nCoverage: The primary limitation is the 87% data void, concentrated in forests, complex vegetation, and water.\n\nUtility: Best suited for urban and open agricultural areas; less reliable in vegetated or complex terrain."
    },
    {
        "title": "Future Work",
        "content": "• Investigate data fusion methods to intelligently fill data voids using other sensors.\n• Test the impact of different InSAR processing parameters on improving coverage.\n• Analyze the influence of seasonality and temporal changes on L-band coherence."
    }
]


# --- Presentation Generation Script ---
def create_presentation(slides_content):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]

    # Define Theme Colors
    TITLE_BLUE = RGBColor(1, 48, 97)
    BODY_GREY = RGBColor(70, 70, 70)

    for i, slide_info in enumerate(slides_content):
        slide = prs.slides.add_slide(blank_slide_layout)

        # Add Title
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1.0))
        p_title = title_shape.text_frame.paragraphs[0]
        p_title.text = slide_info['title']
        p_title.font.name = 'Calibri'
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = TITLE_BLUE

        # Add Separator Line
        line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(0.5), Inches(1.1), Inches(15), Inches(0))
        line.line.color.rgb = TITLE_BLUE
        line.line.width = Pt(2.5)

        # Add Body Content
        body_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(7.5), Inches(6.8))
        body_frame = body_shape.text_frame
        body_frame.word_wrap = True
        p_body = body_frame.paragraphs[0]
        p_body.text = slide_info['content']
        p_body.font.name = 'Calibri'
        p_body.font.size = Pt(20)
        p_body.font.color.rgb = BODY_GREY
        p_body.line_spacing = 1.3

        # Add Image Placeholder on the right side
        if "[Placeholder" in slide_info['content']:
            # This shape is a visual guide in the PPT for where to place the image
            placeholder_shape = slide.shapes.add_textbox(Inches(8.5), Inches(1.4), Inches(7), Inches(6.8))
            p_placeholder = placeholder_shape.text_frame.paragraphs[0]
            p_placeholder.text = "Drop Image/Figure Here"
            p_placeholder.font.size = Pt(24)
            p_placeholder.font.italic = True

        # Add Footer
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(8.5), Inches(15), Inches(0.4))
        p_footer = footer.text_frame.paragraphs[0]
        p_footer.text = f"SAOCOM InSAR Validation Study | Slide {i + 1}"
        p_footer.font.name = 'Calibri'
        p_footer.font.size = Pt(12)
        p_footer.font.color.rgb = BODY_GREY

    output_filename = "Visually_Driven_SAOCOM_Presentation.pptx"
    prs.save(output_filename)
    print(f"✅ Visually-driven presentation successfully generated: '{output_filename}'")


# Run the function
create_presentation(slides_data)